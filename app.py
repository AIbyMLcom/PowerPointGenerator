from flask import Flask, render_template, request, send_from_directory, redirect, jsonify, send_file
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import os
import openai
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import random 
import re

client = openai.OpenAI(api_key='')
app = Flask(__name__)

limiter = Limiter(
    app,
    default_limits=["10 per day"], #This is the rate limit, you can remove it if you want
)

Prompt = """Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has a summary.
-At least 8 slides.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END"""

def create_ppt_text(Input):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": (Prompt)},
            {"role": "user", "content": ("The user wants a presentation about " + Input)}
        ],
        temperature=0.5,
    )

    return response.choices[0].message.content

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{request.host_url}{file_path}"

@app.route('/GeneratedPresentations/<path:path>')
def send_generated_image(path):
    return send_file(f'GeneratedPresentations/{path}', as_attachment=True)
    
@app.route("/powerpoint")
def powerpoint():
    return render_template("powerpoint.html", charset="utf-8")
    
@app.route("/")
def home():
    return render_template("powerpoint.html", charset="utf-8")

@app.route("/get")
@limiter.limit("10 per day, key_func=get_remote_address")
def get_bot_response():
    user_text = request.args.get("msg")
    last_char = user_text[-1]
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    number = 1

    if last_char.isdigit():
        number = int(last_char)
        input_string = user_text[:-2]
        print("Design Number:", number, "selected.")
    else:
        print("No design specified, using default design...")
        
    if number > 7:
        number = 1
        print("Unavailable design, using default design...")
    elif number == 0:
        number = 1
        print("Unavailable design, using default design...")

    # Generate a filename using OpenAI API
    filename_prompt = f"Generate a short, descriptive filename based on the following input: \"{input_string}\". Answer just with the short filename, no other explainment."
    filename_response = client.chat.completions.create(
        model="gpt-3.5-turbo-1106",
        messages=[
            {"role": "system", "content": filename_prompt},
        ],
        temperature=0.5,
        max_tokens=30,
    )
    filename = filename_response.choices[0].message.content.strip().replace(" ", "_")

    with open(f'Cache/{filename}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(input_string))

    pptlink = create_ppt(f'Cache/{filename}.txt', number, filename)
    return str(pptlink)

if __name__ == '__main__':
    # debug=True in the run() parameter if you want to debug
    app.run()
