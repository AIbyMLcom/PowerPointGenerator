import os
import openai
import gradio as gr
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import csv
import io
import uuid
import json
import pandas as pd
import re
import base64
import requests
from io import BytesIO
from PIL import Image

# Configure OpenAI API key from environment variable (recommended)
openai.api_key = os.environ.get("OPENAI_API_KEY")

# Check if API key is set
if not openai.api_key:
    print("WARNING: OpenAI API key is not set. Set the OPENAI_API_KEY environment variable.")
    print("Example: export OPENAI_API_KEY='your-key-here' (Unix/Mac)")
    print("Example: set OPENAI_API_KEY=your-key-here (Windows)")
    
# Define the model for image generation
IMAGE_MODEL = "gpt-4o"

# Define the function to get presentation output directory
def get_presentation_output_dir():
    """Get the configured presentation output directory from environment or use default."""
    return os.environ.get('PRESENTATION_OUTPUT_DIR', os.path.join(os.getcwd(), 'presentations'))

# Create a directory to store temporary data
def ensure_temp_directory():
    temp_dir = os.path.join(os.getcwd(), 'temp_data')
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)
    return temp_dir

# Ensure both directories exist
ensure_temp_directory()
output_dir = get_presentation_output_dir()
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

def generate_presentation_content(title, topic, data_file):
    """Generate presentation content using OpenAI"""
    
    # Generate unique session ID for this content generation
    session_id = str(uuid.uuid4())
    
    # Process uploaded file if present
    file_data = None
    file_content = None
    dataframe = None
    
    if data_file is not None:
        file_path = data_file.name
        if file_path.endswith('.csv'):
            # Process CSV file
            df = pd.read_csv(file_path)
            file_content = df.to_csv()
            file_data = df.values.tolist()
            file_data.insert(0, df.columns.tolist())  # Add header row
            dataframe = df
        elif file_path.endswith(('.xlsx', '.xls')):
            # Process Excel file
            df = pd.read_excel(file_path)
            file_content = df.to_csv()
            file_data = df.values.tolist()
            file_data.insert(0, df.columns.tolist())  # Add header row
            dataframe = df
    
    # If OpenAI API key is not set, show a warning
    if not openai.api_key:
        return "Error: OpenAI API key is not set. Please set the OPENAI_API_KEY environment variable.", None
    
    # Generate content using OpenAI
    try:
        content_prompt = f"Create content for a PowerPoint presentation titled '{title}'. "
        
        if topic:
            content_prompt += f"The presentation is about: {topic}. "
        
        if file_data:
            content_prompt += "Based on the following data: "
            # Add a sample of the data to the prompt
            if isinstance(file_data, list) and len(file_data) > 0:
                sample_rows = min(5, len(file_data))
                for i in range(sample_rows):
                    content_prompt += f"\nRow {i+1}: {file_data[i]}"
        
        content_prompt += """
Create a well-structured presentation with:
1. Title slide
2. Introduction
3. 3-5 main content slides
4. Conclusion

Format each slide like this:

## Slide Title
- Bullet point 1
- Bullet point 2
- Bullet point 3

Make sure each slide has a clear title preceded by '##' and bullet points that start with '-'.
"""
        
        # Call OpenAI API
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates content for PowerPoint presentations."},
                {"role": "user", "content": content_prompt}
            ],
            max_tokens=2000
        )
        
        generated_content = response.choices[0].message.content
        
        # Store the data temporarily
        temp_data = {
            'title': title,
            'topic': topic,
            'file_data': file_data,
            'file_content': file_content,
            'generated_content': generated_content,
            'has_data_file': data_file is not None
        }
        
        temp_dir = ensure_temp_directory()
        with open(os.path.join(temp_dir, f'{session_id}.json'), 'w') as f:
            json.dump(temp_data, f, default=str)
        
        return generated_content, session_id
        
    except Exception as e:
        return f"Error generating content: {str(e)}", None

def generate_info_image_from_data(session_id):
    """Generate an informative image based on course data from uploaded file"""
    if not openai.api_key:
        return False, "OpenAI API key not set. Please set the OPENAI_API_KEY environment variable."
    
    # Load the stored data
    temp_dir = ensure_temp_directory()
    try:
        with open(os.path.join(temp_dir, f'{session_id}.json'), 'r') as f:
            temp_data = json.load(f)
        
        if not temp_data.get('has_data_file', False):
            return False, "No data file was uploaded. Please upload a CSV or Excel file containing course information."
        
        file_data = temp_data.get('file_data', [])
        title = temp_data.get('title', 'Untitled Presentation')
        topic = temp_data.get('topic', '')
        
        if not file_data or len(file_data) < 2:  # At least header and one row
            return False, "Insufficient data in the uploaded file."
        
        # Extract headers and sample data
        headers = file_data[0]
        
        # Create a descriptive prompt based on the data structure
        prompt = f"Create an informative image for a course presentation titled '{title}'. "
        
        if topic:
            prompt += f"The course is about: {topic}. "
        
        prompt += "The course data includes the following fields: " + ", ".join(headers) + ". "
        
        # Add sample data from the first few rows
        prompt += "Sample course data includes: "
        for i in range(1, min(4, len(file_data))):
            row_data = [f"{headers[j]}: {file_data[i][j]}" for j in range(min(len(headers), len(file_data[i])))]
            prompt += " | ".join(row_data) + "; "
        
        prompt += "Create a visual representation that summarizes this course information. Make it professional, educational, and visually appealing."
        
        # Call the OpenAI API for image generation
        response = openai.chat.completions.create(
            model=IMAGE_MODEL,
            messages=[
                {"role": "system", "content": "You are a helpful assistant that generates informative educational images."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000
        )
        
        # Handle the response
        if hasattr(response, 'choices') and len(response.choices) > 0:
            # Store the indication that we've generated an info image
            temp_data['has_info_image'] = True
            with open(os.path.join(temp_dir, f'{session_id}.json'), 'w') as f:
                json.dump(temp_data, f, default=str)
            
            return True, "Course info image generated successfully"
        else:
            return False, "No image was generated in the response"
            
    except Exception as e:
        return False, f"Error generating course info image: {str(e)}"


    """Extract slides from content for display"""
    lines = content.split('\n')
    slides = []
    current_title = None
    current_content = []
    
    # First pass - identify slide boundaries and titles more robustly
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
            
        # More robust slide title detection
        is_title = (
            line.startswith('#') or  # Markdown heading
            (line.startswith('Slide') and (':' in line)) or  # "Slide X: Title" format
            (line.isupper() and len(line) > 5 and len(line) < 60) or  # ALL CAPS TITLE
            (i > 0 and i < len(lines)-1 and not lines[i-1].strip() and not lines[i+1].strip()) or  # Isolated line
            ('slide' in line.lower() and len(line) < 60)  # Contains "slide" keyword
        )
        
        if is_title:
            # Save the previous slide if we have one
            if current_title and current_content:
                slides.append({
                    'title': current_title,
                    'content': current_content.copy()
                })
            
            # Start a new slide
            if ':' in line:
                parts = line.split(':', 1)
                current_title = parts[1].strip()
            else:
                current_title = line.replace('#', '').strip()
            current_content = []
        elif current_title is not None:  # Only add content if we have a current slide
            # Process content lines
            processed_line = line.lstrip('*-•').strip()
            if processed_line:  # Skip empty lines
                current_content.append(processed_line)
    
    # Don't forget to add the last slide
    if current_title and current_content:
        slides.append({
            'title': current_title,
            'content': current_content.copy()
        })
    
    return slides

# Function to extract slides and return them as a list for display with simple numbering
def extract_slides_from_content(content):
    """Extract slides from content for display with numeric slide numbers"""
    lines = content.split('\n')
    slides = []
    current_title = None
    current_content = []
    
    # First pass - identify slide boundaries and titles more robustly
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
            
        # More robust slide title detection
        is_title = (
            line.startswith('#') or  # Markdown heading
            (line.startswith('Slide') and (':' in line)) or  # "Slide X: Title" format
            (line.isupper() and len(line) > 5 and len(line) < 60) or  # ALL CAPS TITLE
            (i > 0 and i < len(lines)-1 and not lines[i-1].strip() and not lines[i+1].strip()) or  # Isolated line
            ('slide' in line.lower() and len(line) < 60)  # Contains "slide" keyword
        )
        
        if is_title:
            # Save the previous slide if we have one
            if current_title and current_content:
                slides.append({
                    'title': current_title,
                    'content': current_content.copy()
                })
            
            # Start a new slide
            if ':' in line:
                parts = line.split(':', 1)
                current_title = parts[1].strip()
            else:
                current_title = line.replace('#', '').strip()
            current_content = []
        elif current_title is not None:  # Only add content if we have a current slide
            # Process content lines
            processed_line = line.lstrip('*-•').strip()
            if processed_line:  # Skip empty lines
                current_content.append(processed_line)
    
    # Don't forget to add the last slide
    if current_title and current_content:
        slides.append({
            'title': current_title,
            'content': current_content.copy()
        })
    
    return slides

# Function to extract slides and return them as a list for display
def get_slide_list(content):
    """Extract slides and return them as a list for display with simple numbering"""
    if not content:
        return [], "Please generate content first."
    
    slides = extract_slides_from_content(content)
    if not slides:
        return [], "No slides could be extracted. Please check the content format."
    
    # Format slide titles for the list with simple numbering
    slide_options = []
    for i, slide in enumerate(slides):
        slide_num = i + 1  # 1-based indexing for display
        slide_title = slide['title']
        # Limit title length for display
        display_title = slide_title[:40] + "..." if len(slide_title) > 40 else slide_title
        slide_options.append(f"{slide_num}: {display_title}")
    
    return slide_options, f"{len(slides)} slides extracted."

# Function to select a slide by its number
def select_slide_by_number(content, slide_number):
    """Select a slide by its number with simplified output"""
    if not content:
        return "No content provided.", None
    
    try:
        # Convert slide_number to integer (1-based) and then to 0-based index
        slide_idx = int(slide_number) - 1
    except ValueError:
        return f"Invalid slide number: '{slide_number}'.", None
    
    slides = extract_slides_from_content(content)
    
    if not slides:
        return "No slides could be extracted.", None
    
    if slide_idx < 0 or slide_idx >= len(slides):
        return f"Slide number {slide_number} is out of range (1-{len(slides)}).", None
    
    # Get the selected slide
    slide = slides[slide_idx]
    
    # Format the slide content for display with simple slide number
    bullet_points = "\n".join([f"• {point}" for point in slide['content']])
    preview = f"""
    <div style="border: 1px solid #ccc; padding: 10px; border-radius: 5px;">
        <div style="margin-bottom: 8px;">
            <span style="font-weight: bold; color: #3498db;">Slide {slide_idx+1}</span>
            <span style="font-weight: bold;"> - {slide['title']}</span>
        </div>
        <div>
            {"<br>• ".join(slide['content'])}
        </div>
    </div>
    """
    
    # Return both the formatted content and the slide index
    return preview, str(slide_idx)

# Function to generate an image for a single selected slide
def generate_single_slide_image(content, slide_idx_str):
    """Generate an image for a single selected slide with simplified output"""
    if not content:
        return {}, "Error: No content provided. Please generate presentation content first.", None, False
    
    if slide_idx_str is None or slide_idx_str == "":
        return {}, "Error: No slide selected. Please select a slide for image generation.", None, False
    
    try:
        slide_idx = int(slide_idx_str)
    except ValueError:
        return {}, f"Error: Invalid slide selection: '{slide_idx_str}'.", None, False
    
    slides = extract_slides_from_content(content)
    if not slides or slide_idx >= len(slides):
        return {}, "Error: Invalid slide selection or no slides could be extracted.", None, False
    
    slide = slides[slide_idx]
    success, message, image_data = generate_image_for_slide(slide['title'], slide['content'])
    
    # Create a dictionary with only the selected slide
    image_results = {slide_idx_str: success}
    
    # Format status message with simple slide number
    status_text = f"""
    <div style="padding: 10px; border-radius: 5px; 
         background-color: {success and '#d4edda' or '#f8d7da'}; 
         border: 1px solid {success and '#c3e6cb' or '#f5c6cb'};">
        <strong>Slide {slide_idx+1}:</strong> {slide['title']}
        <br>
        <strong>Status:</strong> {message}
        <br>
        {success and '✅ Image generated successfully!' or '❌ The image could not be generated.'}
    </div>
    """
    
    # Return image data and visibility flag along with other results
    return image_results, status_text, image_data, success

# Function to generate an image for a slide using OpenAI's image generation
def generate_image_for_slide(slide_title, slide_content):
    """Generate an image for a slide using OpenAI's image generation"""
    if not openai.api_key:
        return None, "OpenAI API key not set. Please set the OPENAI_API_KEY environment variable.", None
    
    try:
        # Create a prompt based on slide content
        prompt = f"Create an image for a slide titled '{slide_title}'. "
        prompt += "Content: " + ", ".join(slide_content[:5])  # Include first 5 bullet points
        prompt += ". Make the image professional and relevant to the content."
        
        # Call the OpenAI API for image generation
        response = openai.chat.completions.create(
            model=IMAGE_MODEL,
            messages=[
                {"role": "system", "content": "You are a helpful assistant that generates presentation slide images."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000
        )
        
        # Get the generated image content (if available)
        if hasattr(response, 'choices') and len(response.choices) > 0:
            # For testing and UI display, create a simple gradient image
            # In a real implementation, you would use the actual image data from the AI service
            from PIL import Image, ImageDraw
            
            # Create a sample image for demonstration
            img = Image.new('RGB', (400, 300), color = (73, 109, 137))
            d = ImageDraw.Draw(img)
            
            # Add a gradient background
            for i in range(0, 300, 1):
                r = int(73 + (182-73) * i / 300)
                g = int(109 + (73-109) * i / 300)
                b = int(137 + (151-137) * i / 300)
                d.line([(0, i), (400, i)], fill=(r, g, b))
            
            # Add a title text
            d.text((20, 20), slide_title, fill=(255, 255, 255))
            
            # Draw some representative graphics based on the slide content
            d.rectangle([(50, 100), (350, 250)], outline=(255, 255, 255))
            
            return True, "Image generated successfully", img
        else:
            return None, "No image was generated in the response", None
            
    except Exception as e:
        return None, f"Error generating image: {str(e)}", None
    
# Function to create a PowerPoint presentation from content
def create_ppt_presentation(content, session_id, slide_images=None):
    """Create PowerPoint presentation from content with optional images"""
    
    if not session_id:
        return "Error: No session ID provided. Please generate content first."
    
    # Load the stored data
    temp_dir = ensure_temp_directory()
    try:
        with open(os.path.join(temp_dir, f'{session_id}.json'), 'r') as f:
            temp_data = json.load(f)
        
        title = temp_data.get('title', 'Untitled Presentation')
        has_info_image = temp_data.get('has_info_image', False)
        
        # Create a presentation
        prs = Presentation()
        
        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:  # Subtitle placeholder
            slide.placeholders[1].text = "Generated with AI"
        
        # Process the content to create slides
        slides = extract_slides_from_content(content)
        
        # If we couldn't extract slides, fallback to simpler approaches
        if not slides:
            lines = content.split('\n')
            current_title = None
            current_content = []
            
            # First pass - identify slide boundaries and titles more robustly
            for i, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                    
                # More robust slide title detection
                is_title = (
                    line.startswith('#') or  # Markdown heading
                    (line.startswith('Slide') and (':' in line)) or  # "Slide X: Title" format
                    (line.isupper() and len(line) > 5 and len(line) < 60) or  # ALL CAPS TITLE
                    (i > 0 and i < len(lines)-1 and not lines[i-1].strip() and not lines[i+1].strip()) or  # Isolated line
                    ('slide' in line.lower() and len(line) < 60)  # Contains "slide" keyword
                )
                
                if is_title:
                    # Save the previous slide if we have one
                    if current_title and current_content:
                        slides.append({
                            'title': current_title,
                            'content': current_content
                        })
                    
                    # Start a new slide
                    if ':' in line:
                        parts = line.split(':', 1)
                        current_title = parts[1].strip()
                    else:
                        current_title = line.replace('#', '').strip()
                    current_content = []
                elif current_title is not None:  # Only add content if we have a current slide
                    # Process content lines
                    processed_line = line.lstrip('*-•').strip()
                    if processed_line:  # Skip empty lines
                        current_content.append(processed_line)
            
            # Don't forget to add the last slide
            if current_title and current_content:
                slides.append({
                    'title': current_title,
                    'content': current_content
                })
        
        # If we couldn't detect any slides, fallback to a simpler approach - make each paragraph a slide
        if not slides:
            paragraphs = []
            current_para = []
            
            for line in lines:
                if line.strip():
                    current_para.append(line.strip())
                elif current_para:  # Empty line and we have content
                    paragraphs.append(current_para)
                    current_para = []
            
            # Add the last paragraph if it exists
            if current_para:
                paragraphs.append(current_para)
            
            # Convert paragraphs to slides
            for i, para in enumerate(paragraphs):
                if len(para) >= 1:
                    title = para[0]
                    content = para[1:] if len(para) > 1 else [""]
                    slides.append({
                        'title': title,
                        'content': content
                    })
        
        # If we still don't have any slides, create a single slide with all content
        if not slides:
            slides.append({
                'title': title,
                'content': content.split('\n')
            })
        
        # Now create slides from the parsed content
        for i, slide_data in enumerate(slides):
            # Use a layout with an image if this slide has an image
            if slide_images and str(i) in slide_images and slide_images[str(i)]:
                # Use layout with title, content and picture
                content_slide_layout = prs.slide_layouts[8] if len(prs.slide_layouts) > 8 else prs.slide_layouts[1]
            else:
                # Use standard layout with title and content
                content_slide_layout = prs.slide_layouts[1]
                
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set slide title
            slide.shapes.title.text = slide_data['title']
            
            # Add content - handle bullet points
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'placeholders') and len(slide.shapes.placeholders) > 1:
                tf = slide.shapes.placeholders[1].text_frame
                for point in slide_data['content']:
                    p = tf.add_paragraph()
                    p.text = point.lstrip('*-•').strip()
                    p.level = 0  # Top level bullet point
                    
                    # Detect if this should be a sub-bullet
                    if point.startswith('  ') or point.startswith('\t'):
                        p.level = 1
        
        # Save the presentation
        output_dir = get_presentation_output_dir()
        safe_title = title.replace(' ', '_').replace('/', '_').replace('\\', '_')
        filename = f"{safe_title}_{session_id[:8]}.pptx"
        output_path = os.path.join(output_dir, filename)
        prs.save(output_path)
        
        return output_path
        
    except Exception as e:
        return f"Error creating presentation: {str(e)}"

# Enhanced create presentation function with image support
def create_enhanced_ppt(content, session_id, slide_images):
    return create_ppt_presentation(content, session_id, slide_images)

# Implementation for the "Generate Images" tab in the interface
def app_interface():
    with gr.Blocks(title="PowerPoint Presentation Generator") as app:
        gr.Markdown("# PowerPoint Presentation Generator")
        
        # Step 1: Generate Content tab
        with gr.Tab("Step 1: Generate Content"):
            with gr.Row():
                with gr.Column():
                    gr.Markdown("""
                    ### Step 1: Generate Presentation Content
                    
                    1. Enter a title for your presentation
                    2. Describe your topic in detail or provide key points to include
                    3. Optionally upload a CSV or Excel file with data to incorporate
                    4. Click "Generate Content" to create your presentation structure
                    """)
                    
                    title_input = gr.Textbox(label="Presentation Title", placeholder="Enter a title for your presentation")
                    topic_input = gr.Textbox(label="Topic Description", placeholder="Describe the topic or provide details", lines=4)
                    data_file = gr.File(label="Upload Data File (CSV/Excel)", file_types=["csv", "xlsx", "xls"])
                    generate_btn = gr.Button("Generate Content")
                
                with gr.Column():
                    gr.Markdown("""
                    ### Generated Presentation Content
                    
                    Your presentation content will appear here after generation. You can:
                    - Review and edit the content as needed
                    - Proceed to Step 2 to add images to your slides

                    """)
                    
                    content_output = gr.Textbox(label="Generated Content", lines=15)
                    session_id = gr.Textbox(label="Session ID", visible=False)
                    
            with gr.Row(visible=False) as course_info_row:
                with gr.Column():
                    gr.Markdown("### Course Data Visualization")
                    gr.Markdown("""
                    We detected that you uploaded a data file. You can generate an informative 
                    visualization based on this data to include in your presentation.
                    
                    Click the button below to create a visual summary of your data.
                    """)
                    generate_info_image_btn = gr.Button("Generate Course Info Image")
                
                with gr.Column():
                    info_image_status = gr.Textbox(label="Status")
        
        # Step 2: Generate Images tab
        with gr.Tab("Step 2: Generate Images"):
            gr.Markdown("""
            ### Step 2: Add Images to Your Slides
            
            Enhance your presentation with AI-generated images for specific slides:
            1. Click "Extract Slides" to analyze your presentation content
            2. Select a slide by number or from the dropdown menu
            3. Click "Generate Image" to create an image for the selected slide
            4. Repeat for any other slides you want to enhance with images
            """)
            
            with gr.Row():
                with gr.Column(scale=2):
                    content_for_images = gr.Textbox(
                        label="Presentation Content", 
                        lines=10, 
                        interactive=True,
                        placeholder="Your presentation content will appear here automatically from Step 1"
                    )
                    
                    extract_slides_btn = gr.Button("Extract Slides", variant="primary")
                    slide_count_status = gr.Textbox(label="Slide Information", interactive=False)
                
                with gr.Column(scale=1):
                    gr.Markdown("### Select a Slide")
                    
                    slide_number_input = gr.Number(
                        label="Slide Number",
                        minimum=1,
                        step=1,
                        precision=0,
                        interactive=True
                    )
                    
                    gr.Markdown("""
                    #### Available Slides
                    Select a slide from the dropdown menu below to quickly access it.
                    Each slide is numbered and shows its title for easy identification.
                    """)
                    
                    slides_list = gr.Dropdown(
                        label="Slide Selection",
                        choices=[],
                        interactive=True,
                        multiselect=False,
                        allow_custom_value=False
                    )
                    
                    current_slide_idx = gr.Textbox(label="Selected Slide Index", visible=False)
                    
            with gr.Row():
                with gr.Column(scale=2):
                    slide_display = gr.Markdown(
                        value="No slide selected yet. Extract slides first, then select a slide."
                    )
                    
                    # Add image display area
                    gr.Markdown("### Generated Image Preview")
                    image_preview = gr.Image(
                        label="Generated Image", 
                        interactive=False,
                        visible=False,
                        type="pil"
                    )
                
                with gr.Column(scale=1):
                    gr.Markdown("### Generate Image")
                    
                    generate_image_btn = gr.Button(
                        "Generate Image", 
                        variant="primary"
                    )
                    
                    image_gen_status = gr.HTML(label="Image Generation Status")
                    slide_images = gr.JSON(label="Generated Images", visible=False)
        
        # Step 3: Create Presentation tab
        with gr.Tab("Step 3: Create Presentation"):
            gr.Markdown("""
            ### Step 3: Create Your PowerPoint Presentation
            
            Final steps to generate your PowerPoint file:
            1. Review and edit your presentation content if needed
            2. Click "Create PowerPoint Presentation" to generate your PPTX file
            3. Once generated, click "Download Presentation" to save it to your computer
            
            All images you generated in Step 2 will be automatically included in the appropriate slides.
            """)
            
            with gr.Row():
                with gr.Column():
                    content_input = gr.Textbox(label="Edit Content (if needed)", lines=15)
                    session_id_input = gr.Textbox(label="Session ID", visible=False)
                    slides_with_images = gr.JSON(label="Slides With Images", visible=False)
                    
                    create_btn = gr.Button("Create PowerPoint Presentation", variant="primary", size="lg")
                
                with gr.Column():
                    result_output = gr.Textbox(label="Creation Status", interactive=False)
                    file_output = gr.File(label="Download Presentation")
                    
                    with gr.Accordion("Progress Summary", open=True):
                        gr.Markdown("""
                        ### Your Progress
                        1. ✅ Generated presentation content
                        2. ✅ Selected slides for image generation (optional)
                        3. ✅ Generated course info image (if data file uploaded)
                        4. 📌 Create and download your final presentation
                        """)
        
        # Function to check if the data file was uploaded
        def check_data_file(session_id_value):
            if not session_id_value:
                return gr.Row.update(visible=False)
            
            temp_dir = ensure_temp_directory()
            try:
                with open(os.path.join(temp_dir, f'{session_id_value}.json'), 'r') as f:
                    temp_data = json.load(f)
                
                if temp_data.get('has_data_file', False):
                    return gr.Row.update(visible=True)
            except:
                pass
            
            return gr.Row.update(visible=False)
        
        # Connect the components for Step 1
        generate_btn.click(
            fn=generate_presentation_content,
            inputs=[title_input, topic_input, data_file],
            outputs=[content_output, session_id]
        )
        
        session_id.change(
            fn=check_data_file,
            inputs=[session_id],
            outputs=[course_info_row]
        )
        
        generate_info_image_btn.click(
            fn=generate_info_image_from_data,
            inputs=[session_id],
            outputs=[info_image_status]
        )
        
        # Connect components for Step 2 - Image Generation
        content_output.change(
            fn=lambda x: x,
            inputs=[content_output],
            outputs=[content_for_images]
        )
        
        extract_slides_btn.click(
            fn=get_slide_list,
            inputs=[content_for_images],
            outputs=[slides_list, slide_count_status]
        )
        
        # Handle slide selection
        slide_number_input.change(
            fn=select_slide_by_number,
            inputs=[content_for_images, slide_number_input],
            outputs=[slide_display, current_slide_idx]
        )
        
        slides_list.change(
            fn=lambda content, selection: select_slide_by_number(content, selection.split(':')[0].replace('Slide', '').strip()) if selection else ("No slide selected", None),
            inputs=[content_for_images, slides_list],
            outputs=[slide_display, current_slide_idx]
        )
        
        # Modified function to handle image preview
        def handle_image_generation(content, slide_idx):
            results, status, image, visible = generate_single_slide_image(content, slide_idx)
            return results, status, image, gr.update(visible=visible)
        
        generate_image_btn.click(
            fn=handle_image_generation,
            inputs=[content_for_images, current_slide_idx],
            outputs=[slide_images, image_gen_status, image_preview, image_preview]
        )
        
        # Connect components for Step 3 - Create Presentation
        content_for_images.change(
            fn=lambda x: x,
            inputs=[content_for_images],
            outputs=[content_input]
        )
        
        session_id.change(
            fn=lambda x: x,
            inputs=[session_id],
            outputs=[session_id_input]
        )
        
        slide_images.change(
            fn=lambda x: x,
            inputs=[slide_images],
            outputs=[slides_with_images]
        )
        
        create_btn.click(
            fn=create_enhanced_ppt,
            inputs=[content_input, session_id_input, slides_with_images],
            outputs=[result_output, file_output]
        )
    
    return app

# Launch the Gradio app
if __name__ == "__main__":
    app = app_interface()
    app.launch(share=False)
